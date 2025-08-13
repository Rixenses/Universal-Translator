"""
Universal Translator
by Ariq S. F. 'Rixenses' Ibrahim

Supported formats / Format yang didukung:
- XML
- TXT
- JSON
- PDF
- DOCX
- XLSX
- PPTX

Dependencies / Dependensi:
pip install googletrans==4.0.0-rc1
pip install python-docx openpyxl python-pptx PyMuPDF
"""

import os
import re
import time
import json
import shutil
import argparse
import fitz  # PyMuPDF
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from googletrans import Translator

# Pola regex untuk mendeteksi komentar XML / Regex pattern to detect XML comments
COMMENT_PATTERN = re.compile(r"<!--\s*EN:.*?-->", re.DOTALL)
# Pola regex untuk menangkap teks di antara tag XML / Regex pattern to capture text between XML tags
TAG_TEXT_PATTERN = re.compile(r">(.*?)<", re.DOTALL)

# Inisialisasi translator Google / Initialize Google translator
translator = Translator()

# =====================
# Fungsi utilitas / Utility functions
# =====================
def chunk_list(lst, size):
    # Membagi list menjadi beberapa bagian kecil / Split list into smaller chunks
    for i in range(0, len(lst), size):
        yield lst[i:i+size]

def translate_batch(texts, target_lang):
    # Menerjemahkan sekumpulan teks sekaligus / Translate a batch of texts at once
    if not texts:
        return []
    try:
        results = translator.translate(texts, dest=target_lang)
        return [res.text for res in results]
    except Exception as e:
        print(f"⚠ Error batch translate: {e}")
        return texts

def backup_file(file_path):
    # Membuat file cadangan sebelum diterjemahkan / Create a backup file before translating
    backup_path = file_path + ".bak"
    shutil.copy(file_path, backup_path)
    print(f"📦 Backup dibuat: {backup_path} / Backup created: {backup_path}")

# =====================
# Handler tiap format / Handlers for each format
# =====================
def process_xml(file_path, target_lang, batch_size):
    # Memproses file XML / Process XML file
    with open(file_path, "r", encoding="utf-8") as f:
        content = f.read()

    parts = re.split(f"({COMMENT_PATTERN.pattern})", content)
    new_parts = []

    for part in parts:
        if COMMENT_PATTERN.fullmatch(part):
            # Jangan terjemahkan komentar / Do not translate comments
            new_parts.append(part)
        else:
            matches = list(TAG_TEXT_PATTERN.finditer(part))
            texts_to_translate = [m.group(1).strip() if m.group(1).strip() else "" for m in matches]

            translated_texts = []
            for batch in chunk_list(texts_to_translate, batch_size):
                translated_texts.extend(translate_batch(batch, target_lang))
                time.sleep(1)  # Jeda agar tidak kena limit / Delay to avoid rate limit

            new_part = part
            for m, translated in zip(matches, translated_texts):
                if translated.strip():
                    new_part = new_part.replace(f">{m.group(1)}<", f">{translated}<", 1)
            new_parts.append(new_part)

    with open(file_path, "w", encoding="utf-8") as f:
        f.write("".join(new_parts))

def process_txt(file_path, target_lang, batch_size):
    # Memproses file TXT / Process TXT file
    with open(file_path, "r", encoding="utf-8") as f:
        lines = f.readlines()

    new_lines = []
    for batch in chunk_list([l.strip("\n") for l in lines], batch_size):
        translated = translate_batch(batch, target_lang)
        new_lines.extend(translated)
        time.sleep(1)

    with open(file_path, "w", encoding="utf-8") as f:
        f.write("\n".join(new_lines))

def process_json(file_path, target_lang, batch_size):
    # Memproses file JSON / Process JSON file
    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    def translate_json(data):
        # Fungsi rekursif untuk menerjemahkan semua teks di JSON / Recursive function to translate all text in JSON
        if isinstance(data, dict):
            return {k: translate_json(v) for k, v in data.items()}
        elif isinstance(data, list):
            return [translate_json(v) for v in data]
        elif isinstance(data, str):
            return translate_batch([data], target_lang)[0]
        return data

    new_data = translate_json(data)
    with open(file_path, "w", encoding="utf-8") as f:
        json.dump(new_data, f, ensure_ascii=False, indent=4)

def process_pdf(file_path, target_lang, batch_size):
    # Memproses file PDF / Process PDF file
    doc = fitz.open(file_path)
    for page in doc:
        text = page.get_text("text")
        lines = text.split("\n")
        translated_lines = []
        for batch in chunk_list(lines, batch_size):
            translated_lines.extend(translate_batch(batch, target_lang))
            time.sleep(1)
        page.clean_contents()
        page.insert_text((50, 50), "\n".join(translated_lines))
    doc.save(file_path)

def process_docx(file_path, target_lang, batch_size):
    # Memproses file DOCX / Process DOCX file
    doc = Document(file_path)
    for para in doc.paragraphs:
        if para.text.strip():
            para.text = translate_batch([para.text], target_lang)[0]
            time.sleep(0.5)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    cell.text = translate_batch([cell.text], target_lang)[0]
                    time.sleep(0.5)
    doc.save(file_path)

def process_xlsx(file_path, target_lang, batch_size):
    # Memproses file XLSX / Process XLSX file
    wb = load_workbook(file_path)
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and cell.value.strip():
                    cell.value = translate_batch([cell.value], target_lang)[0]
                    time.sleep(0.5)
    wb.save(file_path)

def process_pptx(file_path, target_lang, batch_size):
    # Memproses file PPTX / Process PPTX file
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if text.strip():
                    shape.text_frame.text = translate_batch([text], target_lang)[0]
                    time.sleep(0.5)
    prs.save(file_path)

# =====================
# Fungsi utama / Main functions
# =====================
def process_file(file_path, target_lang, batch_size):
    # Memilih handler sesuai tipe file / Select handler based on file type
    ext = os.path.splitext(file_path)[1].lower()
    backup_file(file_path)

    try:
        if ext == ".xml":
            process_xml(file_path, target_lang, batch_size)
        elif ext == ".txt":
            process_txt(file_path, target_lang, batch_size)
        elif ext == ".json":
            process_json(file_path, target_lang, batch_size)
        elif ext == ".pdf":
            process_pdf(file_path, target_lang, batch_size)
        elif ext == ".docx":
            process_docx(file_path, target_lang, batch_size)
        elif ext == ".xlsx":
            process_xlsx(file_path, target_lang, batch_size)
        elif ext == ".pptx":
            process_pptx(file_path, target_lang, batch_size)
        else:
            print(f"⏩ Skipped: {file_path} (unsupported format)")
            return
        print(f"✅ Success: {file_path}")
    except Exception as e:
        print(f"⚠ Failed to process {file_path}: {e}")

def main():
    # Parser argumen CLI / CLI argument parser
    parser = argparse.ArgumentParser(description="Universal Translator")
    parser.add_argument("--path", required=True, help="Path folder/file target / Target folder/file path")
    parser.add_argument("--lang", default="id", help="Kode bahasa target (default: id) / Target language code (default: id)")
    parser.add_argument("--batch", type=int, default=20, help="Jumlah teks per batch (default: 20) / Number of texts per batch (default: 20)")
    args = parser.parse_args()

    if os.path.isfile(args.path):
        # Jika path adalah file tunggal / If path is a single file
        process_file(args.path, args.lang, args.batch)
    else:
        # Jika path adalah folder, proses semua file di dalamnya / If path is a folder, process all files inside
        for root, _, files in os.walk(args.path):
            for file in files:
                process_file(os.path.join(root, file), args.lang, args.batch)

if __name__ == "__main__":
    main()